
import os
from bs4 import BeautifulSoup, NavigableString
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def add_runs_from_element(p, element, font_name, size, color):
    for child in element.children:
        if isinstance(child, NavigableString):
            run = p.add_run()
            run.text = str(child)
        elif child.name == 'span':
            run = p.add_run()
            run.text = child.get_text()
            if 'highlight' in child.get('class', []):
                run.font.bold = True
                run.font.color.rgb = RGBColor(*hex_to_rgb("#4285F4"))
            elif 'inorganic-text' in child.get('class', []):
                run.font.bold = True
                run.font.color.rgb = RGBColor(*hex_to_rgb("#e74c3c"))
            elif 'organic-text' in child.get('class', []):
                run.font.bold = True
                run.font.color.rgb = RGBColor(*hex_to_rgb("#4285F4"))
        elif child.name == 'ul':
            for li in child.find_all('li', recursive=False):
                li_p = p.text_frame.add_paragraph()
                li_p.text = f"• {li.get_text(strip=True)}"
                li_p.font.name = font_name
                li_p.font.size = Pt(size)
                li_p.font.color.rgb = color
                li_p.level = 1
        elif hasattr(child, 'children'):
             add_runs_from_element(p, child, font_name, size, color)


def create_presentation_from_static_html(html_path, output_path):
    with open(html_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'lxml')

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_slide_layout = prs.slide_layouts[6]
    
    COLORS = {'title': RGBColor(66, 133, 244), 'subtitle': RGBColor(44, 62, 80), 'body': RGBColor(85, 85, 85)}
    FONT = 'Poppins'

    for slide_div in soup.find_all('div', class_='slide'):
        slide = prs.slides.add_slide(blank_slide_layout)
        y_offset = Inches(0.5)

        for el in slide_div.find_all(['h1', 'div', 'img'], recursive=False):
            if not el: continue
            if el.name == 'h1':
                txBox = slide.shapes.add_textbox(Inches(0.5), y_offset, Inches(15), Inches(1))
                p = txBox.text_frame.paragraphs[0]
                p.text = el.text; p.font.name = FONT; p.font.size = Pt(32); p.font.color.rgb = COLORS['title']
                y_offset = Inches(1.5)
            
            elif el.name == 'div' and 'grid' in el.get('class', []):
                for i, card in enumerate(el.find_all(class_='card')):
                    left = Inches(0.5 + i * 8)
                    card_y = y_offset
                    for sub_el in card.contents:
                         if isinstance(sub_el, NavigableString): continue
                         if sub_el.name == 'h2':
                            txBox = slide.shapes.add_textbox(left, card_y, Inches(7.5), Inches(0.7))
                            p = txBox.text_frame.paragraphs[0]
                            p.text = sub_el.text; p.font.name = FONT; p.font.size = Pt(22); p.font.color.rgb = COLORS['subtitle']
                            card_y += Inches(0.8)
                         elif sub_el.name in ['p', 'ul', 'h3']:
                            if sub_el.name == 'h3':
                                txBox = slide.shapes.add_textbox(left, card_y, Inches(7.5), Inches(0.5))
                                p = txBox.text_frame.paragraphs[0]
                                p.text = sub_el.text; p.font.name = FONT; p.font.size = Pt(18); p.font.bold = True
                                card_y += Inches(0.6)
                            else:
                                txBox = slide.shapes.add_textbox(left, card_y, Inches(7.5), Inches(3))
                                p = txBox.text_frame.paragraphs[0]
                                p.text = '' # Clear default text
                                add_runs_from_element(p, sub_el, FONT, 16, COLORS['body'])
                                card_y += Inches(3.2)
                         elif sub_el.name == 'img' and sub_el.has_attr('src'):
                            img_path = os.path.join(os.path.dirname(html_path), sub_el['src'])
                            if os.path.exists(img_path):
                                slide.shapes.add_picture(img_path, left, card_y, width=Inches(7))

            elif el.name == 'img' and el.has_attr('src'):
                img_path = os.path.join(os.path.dirname(html_path), el['src'])
                if os.path.exists(img_path):
                    slide.shapes.add_picture(img_path, Inches(0.5), y_offset, width=Inches(15))

    prs.save(output_path)
    print(f"Presentación generada con éxito en: {output_path}")

if __name__ == '__main__':
    create_presentation_from_static_html('reports/reporte_final_estatico.html', 'reports/reporte_final_estatico.pptx')
